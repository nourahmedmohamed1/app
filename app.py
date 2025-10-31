"""
تطبيق Streamlit لتصدير النماذج
سهل جداً ويشتغل على Streamlit Cloud مجاناً!

للتثبيت المحلي:
pip install streamlit python-pptx reportlab arabic-reshaper python-bidi

للتشغيل:
streamlit run app.py
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_RIGHT
from io import BytesIO
from datetime import datetime

try:
    import arabic_reshaper
    from bidi.algorithm import get_display
    ARABIC_SUPPORT = True
except:
    ARABIC_SUPPORT = False

# إعدادات الصفحة
st.set_page_config(
    page_title="نموذج تصدير البيانات",
    page_icon="📋",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# CSS مخصص
st.markdown("""
<style>
    .main {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    }
    .stApp {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    }
    div[data-testid="stForm"] {
        background: white;
        padding: 2rem;
        border-radius: 20px;
        box-shadow: 0 20px 60px rgba(0, 0, 0, 0.3);
    }
    .css-1d391kg {
        background: white;
    }
    h1, h2, h3 {
        color: #667eea !important;
    }
    .stButton>button {
        width: 100%;
        background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
        color: white;
        font-weight: bold;
        border: none;
        padding: 0.75rem;
        border-radius: 10px;
        font-size: 1.1rem;
    }
    .stButton>button:hover {
        transform: translateY(-2px);
        box-shadow: 0 10px 25px rgba(79, 172, 254, 0.4);
    }
    .success-box {
        background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
        font-weight: bold;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)


def reshape_arabic_text(text):
    """تحويل النص العربي للعرض الصحيح"""
    if not text or not ARABIC_SUPPORT:
        return text
    try:
        reshaped_text = arabic_reshaper.reshape(text)
        bidi_text = get_display(reshaped_text)
        return bidi_text
    except:
        return text


def create_powerpoint(data):
    """إنشاء ملف PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # الشريحة الأولى - عنوان
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = (102, 126, 234)

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "بيانات النموذج"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = (255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    date_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%Y/%m/%d")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(24)
    date_para.font.color.rgb = (255, 255, 255)
    date_para.alignment = PP_ALIGN.CENTER

    # الشريحة الثانية - المعلومات الشخصية
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "المعلومات الشخصية"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(36)
    title_para2.font.bold = True
    title_para2.font.color.rgb = (102, 126, 234)

    content_y = 1.5
    personal_info = [
        f"الاسم الكامل: {data['fullName']}",
        f"البريد الإلكتروني: {data['email']}",
        f"رقم الهاتف: {data['phone']}",
        f"العنوان: {data['address']}"
    ]

    for info in personal_info:
        text_box = slide2.shapes.add_textbox(Inches(1), Inches(content_y), Inches(8), Inches(0.6))
        text_frame = text_box.text_frame
        text_frame.text = info
        para = text_frame.paragraphs[0]
        para.font.size = Pt(24)
        content_y += 0.8

    # الشريحة الثالثة - معلومات العمل
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame3 = title_box3.text_frame
    title_frame3.text = "معلومات العمل"
    title_para3 = title_frame3.paragraphs[0]
    title_para3.font.size = Pt(36)
    title_para3.font.bold = True
    title_para3.font.color.rgb = (102, 126, 234)

    content_y = 1.5
    work_info = [
        f"الشركة: {data['company']}",
        f"المسمى الوظيفي: {data['position']}",
        f"القسم: {data['department']}",
        f"سنوات الخبرة: {data['experience']}"
    ]

    for info in work_info:
        text_box = slide3.shapes.add_textbox(Inches(1), Inches(content_y), Inches(8), Inches(0.6))
        text_frame = text_box.text_frame
        text_frame.text = info
        para = text_frame.paragraphs[0]
        para.font.size = Pt(24)
        content_y += 0.8

    # الشريحة الرابعة - التفاصيل الإضافية
    if data['skills'] or data['notes']:
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame4 = title_box4.text_frame
        title_frame4.text = "تفاصيل إضافية"
        title_para4 = title_frame4.paragraphs[0]
        title_para4.font.size = Pt(36)
        title_para4.font.bold = True
        title_para4.font.color.rgb = (102, 126, 234)

        content_y = 1.5
        if data['skills']:
            text_box = slide4.shapes.add_textbox(Inches(1), Inches(content_y), Inches(8), Inches(0.6))
            text_frame = text_box.text_frame
            text_frame.text = f"المهارات: {data['skills']}"
            para = text_frame.paragraphs[0]
            para.font.size = Pt(20)
            content_y += 1

        if data['notes']:
            text_box = slide4.shapes.add_textbox(Inches(1), Inches(content_y), Inches(8), Inches(2))
            text_frame = text_box.text_frame
            text_frame.text = f"ملاحظات:\n{data['notes']}"
            para = text_frame.paragraphs[0]
            para.font.size = Pt(18)
            text_frame.word_wrap = True

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


def create_pdf(data):
    """إنشاء ملف PDF"""
    pdf_io = BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=letter, rightMargin=72, leftMargin=72)
    elements = []
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'ArabicTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#667eea'),
        spaceAfter=30,
        alignment=TA_RIGHT,
        fontName='Helvetica-Bold'
    )

    heading_style = ParagraphStyle(
        'ArabicHeading',
        parent=styles['Heading2'],
        fontSize=20,
        textColor=colors.HexColor('#667eea'),
        spaceAfter=12,
        spaceBefore=20,
        alignment=TA_RIGHT,
        fontName='Helvetica-Bold'
    )

    elements.append(Paragraph(reshape_arabic_text("بيانات النموذج"), title_style))
    elements.append(Paragraph(f"التاريخ: {datetime.now().strftime('%Y-%m-%d')}", heading_style))
    elements.append(Spacer(1, 20))

    # المعلومات الشخصية
    elements.append(Paragraph(reshape_arabic_text("المعلومات الشخصية"), heading_style))
    personal_data = [
        [reshape_arabic_text('الاسم الكامل'), reshape_arabic_text(data['fullName'])],
        [reshape_arabic_text('البريد الإلكتروني'), data['email']],
        [reshape_arabic_text('رقم الهاتف'), data['phone']],
        [reshape_arabic_text('العنوان'), reshape_arabic_text(data['address'])]
    ]

    personal_table = Table(personal_data, colWidths=[200, 300])
    personal_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#667eea')),
        ('TEXTCOLOR', (0, 0), (0, -1), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'RIGHT'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
        ('TOPPADDING', (0, 0), (-1, -1), 12),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('BACKGROUND', (1, 0), (1, -1), colors.beige)
    ]))
    elements.append(personal_table)
    elements.append(Spacer(1, 20))

    # معلومات العمل
    elements.append(Paragraph(reshape_arabic_text("معلومات العمل"), heading_style))
    work_data = [
        [reshape_arabic_text('الشركة'), reshape_arabic_text(data['company'])],
        [reshape_arabic_text('المسمى الوظيفي'), reshape_arabic_text(data['position'])],
        [reshape_arabic_text('القسم'), reshape_arabic_text(data['department'])],
        [reshape_arabic_text('سنوات الخبرة'), data['experience']]
    ]

    work_table = Table(work_data, colWidths=[200, 300])
    work_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#667eea')),
        ('TEXTCOLOR', (0, 0), (0, -1), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'RIGHT'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
        ('TOPPADDING', (0, 0), (-1, -1), 12),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('BACKGROUND', (1, 0), (1, -1), colors.beige)
    ]))
    elements.append(work_table)
    elements.append(Spacer(1, 20))

    # التفاصيل الإضافية
    if data['skills'] or data['notes']:
        elements.append(Paragraph(reshape_arabic_text("تفاصيل إضافية"), heading_style))
        if data['skills']:
            elements.append(Paragraph(f"{reshape_arabic_text('المهارات')}: {reshape_arabic_text(data['skills'])}", heading_style))
        if data['notes']:
            elements.append(Paragraph(f"{reshape_arabic_text('ملاحظات')}: {reshape_arabic_text(data['notes'])}", heading_style))

    doc.build(elements)
    pdf_io.seek(0)
    return pdf_io


# واجهة التطبيق
st.markdown("<h1 style='text-align: center; color: white;'>📋 نموذج إدخال البيانات</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align: center; color: white; font-size: 1.2rem;'>املأ النموذج وصدّر البيانات كـ PowerPoint أو PDF احترافي</p>", unsafe_allow_html=True)
st.markdown("<br>", unsafe_allow_html=True)

with st.form("data_form"):
    st.markdown("### 👤 المعلومات الشخصية")
    col1, col2 = st.columns(2)

    with col1:
        fullName = st.text_input("الاسم الكامل *", placeholder="أدخل الاسم الكامل")
        phone = st.text_input("رقم الهاتف", placeholder="+20 XXX XXX XXXX")

    with col2:
        email = st.text_input("البريد الإلكتروني *", placeholder="example@email.com")
        address = st.text_input("العنوان", placeholder="المدينة، الدولة")

    st.markdown("### 💼 معلومات العمل")
    col3, col4 = st.columns(2)

    with col3:
        company = st.text_input("اسم الشركة/المؤسسة", placeholder="اسم الشركة")
        department = st.selectbox("القسم", ["", "تقنية المعلومات", "الموارد البشرية", "المالية", "التسويق", "المبيعات", "خدمة العملاء", "الإدارة", "أخرى"])

    with col4:
        position = st.text_input("المسمى الوظيفي", placeholder="المسمى الوظيفي")
        experience = st.number_input("سنوات الخبرة", min_value=0, max_value=50, value=0)

    st.markdown("### 📝 تفاصيل إضافية")
    skills = st.text_area("المهارات", placeholder="أدخل المهارات مفصولة بفواصل")
    notes = st.text_area("ملاحظات", placeholder="أضف أي ملاحظات أو تفاصيل إضافية هنا...")

    st.markdown("<br>", unsafe_allow_html=True)

    col_btn1, col_btn2, col_btn3 = st.columns(3)

    with col_btn1:
        submit_pptx = st.form_submit_button("📊 تصدير PowerPoint", use_container_width=True)
    with col_btn2:
        submit_pdf = st.form_submit_button("📄 تصدير PDF", use_container_width=True)
    with col_btn3:
        st.form_submit_button("🔄 مسح النموذج", use_container_width=True)

# معالجة التصدير
if submit_pptx or submit_pdf:
    if not fullName or not email:
        st.error("⚠️ يرجى ملء الحقول المطلوبة (الاسم الكامل والبريد الإلكتروني)")
    else:
        data = {
            'fullName': fullName,
            'email': email,
            'phone': phone or 'غير محدد',
            'address': address or 'غير محدد',
            'company': company or 'غير محدد',
            'position': position or 'غير محدد',
            'department': department or 'غير محدد',
            'experience': str(experience),
            'skills': skills or 'غير محدد',
            'notes': notes or 'غير محدد'
        }

        with st.spinner('جاري التصدير...'):
            if submit_pptx:
                file_io = create_powerpoint(data)
                filename = f"form_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                st.download_button(
                    label="⬇️ تحميل PowerPoint",
                    data=file_io,
                    file_name=filename,
                    mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
                st.success("✅ تم إنشاء ملف PowerPoint بنجاح!")

            if submit_pdf:
                file_io = create_pdf(data)
                filename = f"form_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                st.download_button(
                    label="⬇️ تحميل PDF",
                    data=file_io,
                    file_name=filename,
                    mime='application/pdf'
                )
                st.success("✅ تم إنشاء ملف PDF بنجاح!")

# معلومات في الأسفل
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("<div style='text-align: center; color: white; padding: 20px;'>© 2025 - تطبيق تصدير البيانات | Powered by Streamlit</div>", unsafe_allow_html=True)